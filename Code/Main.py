# -*- coding: utf-8 -*-
"""
Created on Wed Jun 17 16:05:30 2020

@author: ZuroChang
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

Root=os.path.dirname(os.path.dirname(os.getcwd()))+'\\'
CurrentFolder="PTTAutomation\\"
Project=''

##########################################

CodeFolder=Root+CurrentFolder+"Code\\"+Project
ImportFolder=Root+CurrentFolder+"Import\\"+Project
OutputFolder=Root+CurrentFolder+"Output\\"+Project
TemplateFolder=Root+CurrentFolder+'Template\\'

class ReadExcel:
    '''
    Description:
        Read the excel file including the ppt content. The output is the format
        required to output the ppt file
    '''
    
    Title=[]
    Pictures=[]
    Content=[]
    Template=[]
    SlideLayout=[]
    SlideMapping=[]
    
    def __init__(self,ExcelLocation):
        self.Excel=pd.read_excel(ExcelLocation)\
            .sort_values(by='ID').reset_index(drop=True)
        
    def GetSlideMapping(self):
        MappingRef=pd.read_excel(TemplateFolder+'SlideMapping.xlsx')

        for Template in list(MappingRef['Template'].drop_duplicates()):
            TemplateMap=MappingRef[MappingRef['Template']==Template].reset_index(drop=True)
            
            Mapping=[]
            for i0 in range(TemplateMap.shape[0]):
                Mapping.append(
                    {'Layout':int(TemplateMap['Layout'].iloc[i0])
                    ,'Code':TemplateMap['Code'].iloc[i0]}
                )    
            
            self.SlideMapping.append({'Template':Template,'Mapping':Mapping})
    
    def OutputTitle(self):
        '''
        Description:
            Output self.Title
        '''
        self.Title=list(self.Excel['Title'])
        
    def OutputContent(self):
        '''
        Description:
            Output self.Content
        '''
        def GetBulletLevel(Bullet):
            Bulletlevel={
                0:[str(i) for i in range(1,100)],
                1:[chr(i) for i in range(ord('A'),ord('Z'))],
                2:[chr(i) for i in range(ord('a'),ord('z'))],
            }
            
            for key in Bulletlevel.keys():
                if Bullet in Bulletlevel[key]:
                    return(key)
        
        for i0 in range(self.Excel.shape[0]):
            entry=str(self.Excel['Content'].iloc[i0]).split('\n')
            
            if entry[0]=='nan':
                self.Content.append(None)
            else:
                entryContent=[]
                for i1 in range(len(entry)):
                    Separation=entry[i1].find('. ')
                    Bullet=entry[i1][:Separation]
                    Paragraph=entry[i1][Separation+2:]
                    
                    entryContent.append(
                        {'lvl':GetBulletLevel(Bullet)
                        ,'length':len(Paragraph)
                        ,'Text':Paragraph}
                    )
                
                self.Content.append(entryContent)
        
    def OutputPictures(self):
        '''
        Description:
            Output self.Pictures
        '''
        for i0 in range(self.Excel.shape[0]):
            self.Pictures.append(self.Excel['Pictures'].iloc[i0].split(','))
        
    def OutputSlideLayout(self):
        '''
        Description:
            Output self.Template and self.SlideLayout
        '''
        for i0 in range(self.Excel.shape[0]):
            Template=self.Excel['Template'].iloc[i0]
            for entry in self.SlideMapping:
                if Template==entry['Template']:
                    Mapping=entry['Mapping']
                    break
            
            Code=self.Excel['Slide'].iloc[i0]
            for entry in Mapping:
                if Code==entry['Code']:
                    Layout=entry['Layout']
                    break
                
            self.Template.append(Template)
            
            self.SlideLayout.append(
                {'SlideLayout':Layout
                ,'ContentBlock':int(Code[2])}
            )
    
    
        
    def Run(self):
        '''
        Description:
            Return the content for outputing ppt file.
        '''
        self.GetSlideMapping()
        self.OutputTitle()
        self.OutputPictures()
        self.OutputContent()
        self.OutputSlideLayout()
        
        return(
            pd.DataFrame(
                {'Title':self.Title
                ,'Pictures':self.Pictures
                ,'Content':self.Content
                ,'SlideLayout':self.SlideLayout
                ,'Template':self.Template}
            )
        )

class OutputSlide:
    WordSize={
        'FontSize':[32,28,24,20,18,16,14,12,10],
        'Width':[406400,355600,304800,254000,228600,203200,177800,152400,127000],
        'height':[633747,506997,422498,362141,362141,316873,281665,253498,162455]
    }
    MinTitlePt=18
    
    def __init__(self,Template):
        self.prs = Presentation(TemplateFolder+'{}.pptx'.format(Template))

    def GetPt(self,width,height,WordNum):
        WideCapacity=[width//entry for entry in self.WordSize['Width']]
        
        HeightCapacity=[]
        for entry in self.WordSize['height']:
            if entry:
                HeightCapacity.append(height//entry)
            else:
                HeightCapacity.append(None)
        
        for i0 in range(len(HeightCapacity)):
            if not HeightCapacity[len(HeightCapacity)-i0-1]:
                HeightCapacity[len(HeightCapacity)-i0-1]=HeightCapacity[len(HeightCapacity)-i0]
        
        for i0 in range(len(self.WordSize['FontSize'])):
            if WideCapacity[i0]*HeightCapacity[i0]>=WordNum:
                return(self.WordSize['FontSize'][i0])
    
    def __InsertTitle(self,Shape):
        Shape.text=self.Data['Title']
        Shape.text_frame.paragraphs[0].font.size = \
            Pt(max(self.GetPt(Shape.width,Shape.height,len(self.Data['Title'])),self.MinTitlePt))
    
    def __SeparateContent(self):
        Content=self.Data['Content']
        BlockNum=self.Data['SlideLayout']['ContentBlock']
        if BlockNum:
            Level0=[i0 for i0 in range(len(Content)) if Content[i0]['lvl']==0]
            
            Quotient,Remainder=divmod(len(Level0),BlockNum)
            Separation=[Quotient]*BlockNum
            i0=0
            while Remainder>0:
                Separation[i0]+=1
                Remainder-=1
            
            i0=0
            count=0
            SubContent=[]
            SeparationContent=[]
            for entry in Content:
                if entry['lvl']==0:
                    count+=1
                
                if count<=Separation[i0]:
                    SubContent.append(entry)
                else:
                    SeparationContent.append(SubContent)
                    SubContent=[entry]
                    count=1
                    i0+=1
            else:
                SeparationContent.append(SubContent)
        
            return(SeparationContent)
        else:
            return([])
    
    def __InsertContent(self,Shape,Content):
        text_frame=Shape.text_frame
            
        p=text_frame.paragraphs[0]
        p.text=Content[0]['Text']
        p.level=Content[0]['lvl']
    
        for entry in Content[1:]:
            p=text_frame.add_paragraph()
            p.text=entry['Text']
            p.level=entry['lvl']
    
    def __InsertPicture(self,Shape,Picture):
        self.Slide.shapes.add_picture(Picture,
            Shape.left,Shape.top,Shape.width,Shape.height
        )
    
    def AddSlide(self,Data,PictureFolder):
        self.Slide=self.prs.slides.add_slide(
            self.prs.slide_layouts[Data['SlideLayout']['SlideLayout']]
        )
        self.Data=Data
        
        #B:Get the path of pictures--Pictures     
        Pictures=[]
        for entry in Data['Pictures']:
            for entry1 in os.listdir(PictureFolder):
                if entry in entry1:
                    Pictures.append(PictureFolder+entry1)
                    break
        #E:Get the path of pictures--Pictures    
        
        SeparationContent=self.__SeparateContent()
        count=1
        i0=0
        for shape in self.Slide.placeholders:
            if 'Title'==shape.name[:len('Title')]:
                self.__InsertTitle(shape)
            elif 'Picture'==shape.name[:len('Picture')]:
                if count<=len(self.Data['Pictures']):
                    self.__InsertPicture(shape,Pictures[count-1])
                    count+=1
            elif 'Text'==shape.name[:len('Text')]:
                self.__InsertContent(shape,SeparationContent[i0])
                i0+=1
    
    def Close(self,PPTName='Output',OutputFolder='Output/'):
        self.prs.save('{}{}.pptx'.format(OutputFolder,PPTName))

class OutputPPT:
    ErrorList=[]
    def __init__(self,FolderPath,ExcelName):
        self.FolderPath=FolderPath
        Excel=ReadExcel(FolderPath+ExcelName).Run()
        self.Excel=Excel.reset_index(drop=False).sort_values(by=['Template','index'])
    
    def Run(self,OutputFolder,PPTName,Separation_Flag=0):
        for Template in list(self.Excel['Template'].drop_duplicates()):
            Data=self.Excel[self.Excel['Template']==Template].reset_index(drop=True)

            if Separation_Flag:
                for i0 in range(Data.shape[0]):
                    try:
                        PPT=OutputSlide(Template)
                        PPT.AddSlide(
                            Data=Data.iloc[i0],
                            PictureFolder=self.FolderPath
                        )
                        PPT.Close(
                            PPTName='{}_p{}'.format(PPTName,Data.iloc[i0]['index']),
                            OutputFolder=OutputFolder
                        )
                    except:
                        self.ErrorList.append(Data.iloc[i0]['index'])
            else:
                PPT=OutputSlide(Template)
                for i0 in range(Data.shape[0]):
                    try:
                        PPT.AddSlide(
                            Data=Data.iloc[i0],
                            PictureFolder=self.FolderPath
                        )
                    except:
                        self.ErrorList.append(Data.iloc[i0]['index'])
                PPT.Close(
                    PPTName='{}_{}'.format(PPTName,Template),
                    OutputFolder=OutputFolder
                )
        
        if self.ErrorList:
            for entry in self.ErrorList:
                print('{} has type error!!'.format(entry+1))
        else:
            print('Success!')
            
if __name__=='__main__':       
    FileFolder='ReportGSAM_ZuroChang_20181016/'
    ExcelFile='ReportGSAM_ZuroChang_20181016.xlsx'
    PPTName='Test'
    Separation_Flag=0
    
    A1=OutputPPT(
        FolderPath=ImportFolder+FileFolder,
        ExcelName=ExcelFile
    )
    
    A1.Run(OutputFolder,PPTName,Separation_Flag)

