# -*- coding: utf-8 -*-
"""
Created on Thu Oct  4 09:54:50 2018

@author: gash8
"""

import pandas as pd
import random
import math
from pptx import Presentation
from pptx.util import Pt
import os

#Begin:判斷字體是否為中文
def ContainChese(String):
    if u'\u4e00' <= String <= u'\u9fff':
            return True
    return False
#End:判斷字體是否為中文
#Begin:初始設定
RootPath='C:\\Users\\USER\\Python Script\\PPTAutomation\\'
TemplatePath = RootPath+'Template01.pptx'
ExcelPath = RootPath+'ReportMATrainingII_ZuroChang_20191202_t1.xlsx'
PictureFilePath = RootPath+'ReportMATrainingII_ZuroChang'
OutputName='ReportMATrainingII_ZuroChang_20191202'

WordWidth1 = 123224
WordWidth2 = 113622
WordWidth3 = 100562
WordHeight1 = 321089
WordHeight2 = 285412
WordHeight3 = 265871

Level1 = [str(each)+'. ' for each in range(1,21)]
Level2 = [chr(each)+'. ' for each in range(ord('A'), ord('Z') + 1)]
Level3 = [chr(each)+'. ' for each in range(ord('a'), ord('z') + 1)]
#載入Map檔
Map = pd.read_excel(RootPath+'StyleMap.xlsx')
#End:初始設定
prs = Presentation(TemplatePath) #打開PPT
#Begin:分類模板
PictureInformation = []
TotalLayouts = prs.slide_layouts
for i0 in range(len(TotalLayouts)):
    Slide = prs.slides.add_slide(TotalLayouts[i0]) 
    ShapesName = [ shape.name for shape in Slide.shapes ]
    Title = []
    Content = []
    Picture = []
    for i1 in range(len(ShapesName)):
        if 'Title' in ShapesName[i1]:
            Title.append(int(float(ShapesName[i1][-1]))-1)
        elif 'Picture' in ShapesName[i1]:
            Picture.append(int(float(ShapesName[i1][-1]))-1)
            PictureInfo = pd.DataFrame({'Name':[ShapesName[i1]],
                                        'Top':[Slide.shapes[i1].top],
                                        'Left':[Slide.shapes[i1].left],
                                        'ShapeNum':[int(float(ShapesName[i1][-1]))-1],
                                        'Slide':[i0]})
            if len(PictureInformation)>0:
                PictureInformation = PictureInformation.append(PictureInfo)
            else:
                PictureInformation = PictureInfo
        elif 'Content' in ShapesName[i1]:
            Content.append(int(float(ShapesName[i1][-1]))-1)
    if len(Picture)>0:
        Information = pd.DataFrame({'Type':[i0],
                                    'Title':[Title],
                                    'Content':[Content],
                                    'Picture':[Picture],
                                    'PictureNum':[len(Picture)]})
    else:
        Information = pd.DataFrame({'Type':[i0],
                                    'Title':[Title],
                                    'Content':[Content],
                                    'Picture':[''],
                                    'PictureNum':[0]})
    if i0==0:
        SlideInfo = Information
    else:
        SlideInfo = SlideInfo.append(Information)
#End:分類模板
#Begin:尋找所有圖片檔
PicturePathList = os.listdir(PictureFilePath)
PictureID = [each[0:3] for each in PicturePathList] #分割出ID
PictureName = [each[0:6] for each in PicturePathList] #分割出ID
PicturePath = pd.DataFrame({'ID':PictureID,
                            'Name':PictureName,
                            'Path':PicturePathList}) #建立Map表

#End:尋找所有圖片檔
prs = Presentation(TemplatePath)
Data = pd.read_excel(ExcelPath)
for i0 in range(len(Data)):
    SlideID = Data.ID[i0]   #ID
    SlideTitle = Data.Header[i0]  #標題 
    SlideText = Data.Content[i0]  #內容
    #Begin:建立圖片清單
    PathList = PicturePath.Path[PicturePath.ID.astype(float)==SlideID].tolist()
    PathList = [PictureFilePath+'\\'+each for each in PathList]
    #End:建立圖片清單
    SlidePictureNum = len(PathList) #圖片數量
    #Begin:決定要用哪個模板
    if Data.Template[i0]!=999:
        SlideType = Map.Number[Map.Code==Data.Style[i0]].tolist()[0]
    else:
        RandomList = SlideInfo.Style[SlideInfo.PictureNum==SlidePictureNum].tolist()
        SlideType = random.sample(RandomList,1)[0]            
    #End:決定要用哪個模板
    #Begin:針對內容進行處理
    #計算字數
    if str(SlideText)!='nan':
        TextLen = sum([ContainChese(each)+1 for each in SlideText])
        SlideText = SlideText.split('\n')
        #分段
        TextStop = []
        for i1 in range(len(SlideText)):
            if sum([each in SlideText[i1] for each in Level1])>0:
                TextStop.append(i1)
        TextStop.append(len(SlideText))
    #End:針對內容進行處理
    #Begin:新增PPT
    ContentPlaceholder = SlideInfo.Content[SlideInfo.Type==SlideType][0]
    NewPage = True
    SlideCount = 1
    #如果有1個以上的文字內容區塊
    if len(ContentPlaceholder)>0 and str(SlideText)!='nan':
        #分段插入內容
        for i1 in range(1,len(TextStop)):
            ContentInserted = SlideText[TextStop[i1-1]:TextStop[i1]]
            if NewPage==True:
                TypeSlideLayout = prs.slide_layouts[SlideType]
                Slide = prs.slides.add_slide(TypeSlideLayout)
                Title = Slide.shapes.title
                Title.text_frame.text = Data.Header[i0] 
                TitleLen = 0
                for i2 in range(len(Data.Header[i0])):
                    if ContainChese(Data.Header[i0][i2]):
                        TitleLen = TitleLen+2
                    else:
                        TitleLen = TitleLen+1
                if TitleLen>=90:
                    Title.text_frame.paragraphs[0].font.size = Pt(20)
                else:
                    Title.text_frame.paragraphs[0].font.size = Pt(24)
                NewPage = False
                MaximumCount = len(ContentPlaceholder)
                ContentCount = 0
                ContentHeight = 0
                Add = True
            else:
                Slide = prs.slides[-1]
            #分行插入內容
            InsertedPlaceholder = ContentPlaceholder[ContentCount]
            for i2 in range(len(ContentInserted)):
                #段落層級
                WordNumber = 0
                for i3 in range(len(ContentInserted[i2][3:])):
                    if ContainChese(ContentInserted[i2][3:][i3]):
                        WordNumber = WordNumber+2
                    else:
                        WordNumber = WordNumber+1
                if ContentInserted[i2][:3] in Level1:
                    ParagraphLevel = 0
                    ContentHeight = ContentHeight+math.ceil(WordNumber*WordWidth1/Slide.shapes[InsertedPlaceholder].width)*WordHeight1
                elif ContentInserted[i2][:3] in Level2:
                    ParagraphLevel = 1
                    ContentHeight = ContentHeight+math.ceil(WordNumber*WordWidth2/Slide.shapes[InsertedPlaceholder].width)*WordHeight2
                elif ContentInserted[i2][:3] in Level3:
                    ParagraphLevel = 2
                    ContentHeight = ContentHeight+math.ceil(WordNumber*WordWidth3/Slide.shapes[InsertedPlaceholder].width)*WordHeight3
                if Add == True:
                    NewParagraph = Slide.shapes[InsertedPlaceholder]
                    NewParagraph.text = ContentInserted[i2][3:]
                    Add = False
                else:
                    NewParagraph = Slide.shapes[InsertedPlaceholder].text_frame.add_paragraph()
                    NewParagraph.text = ContentInserted[i2][3:]
                    NewParagraph.level = ParagraphLevel
                
            #換文字區塊
            if i1!=(len(TextStop)-1):
                if ContentHeight>=Slide.shapes[InsertedPlaceholder].height*0.9:
                    if ContentCount<(MaximumCount-1):
                        ContentCount = ContentCount+1
                        Add = True
                    else:
                        NewPage = True
                        SlideCount = SlideCount+1
        #插入圖片
        if len(Data.Location[i0])>0:
            PictureOrder = Data.Location[i0].split(',')
            PicturePlaceholder = SlideInfo.Picture[SlideInfo.Type==SlideType][0]
            if SlideCount==1:
                Slide = prs.slides[-1]
                for i1 in range(len(PictureOrder)):
                    if i1<len(PicturePlaceholder):
                        Pic = PictureFilePath+'\\'+PicturePath.Path[PicturePath.Name==PictureOrder[i1]].tolist()[0]
                        Slide.shapes.add_picture(Pic, Slide.shapes[PicturePlaceholder[i1]].left, Slide.shapes[PicturePlaceholder[i1]].top, Slide.shapes[PicturePlaceholder[i1]].width, Slide.shapes[PicturePlaceholder[i1]].height)
            else:
                for i1 in range(SlideCount):
                    Slide = prs.slides[-1-i1]
                    for i2 in range(len(PictureOrder)):
                        if i2<len(PicturePlaceholder):
                            Pic = PictureFilePath+'\\'+PicturePath.Path[PicturePath.Name==PictureOrder[i2]].tolist()[0]
                            Slide.shapes.add_picture(Pic, Slide.shapes[PicturePlaceholder[i2]].left, Slide.shapes[PicturePlaceholder[i2]].top, Slide.shapes[PicturePlaceholder[i2]].width, Slide.shapes[PicturePlaceholder[i2]].height)
    elif str(SlideText)=='nan':
        TypeSlideLayout = prs.slide_layouts[SlideType]
        Slide = prs.slides.add_slide(TypeSlideLayout)
        Title = Slide.shapes.title
        Title.text_frame.text = Data.Header[i0] 
        TitleLen = 0
        for i2 in range(len(Data.Header[i0])):
            if ContainChese(Data.Header[i0][i2]):
                TitleLen = TitleLen+2
            else:
                TitleLen = TitleLen+1
        if TitleLen>=90:
            Title.text_frame.paragraphs[0].font.size = Pt(20)
        else:
            Title.text_frame.paragraphs[0].font.size = Pt(24)
        PictureOrder = Data.Location[i0].split(',')
        PicturePlaceholder = SlideInfo.Picture[SlideInfo.Type==SlideType][0]
        for i1 in range(len(PictureOrder)):
            if i1<len(PicturePlaceholder):
                Pic = PictureFilePath+'\\'+PicturePath.Path[PicturePath.Name==PictureOrder[i1]].tolist()[0]
                Slide.shapes.add_picture(Pic, Slide.shapes[PicturePlaceholder[i1]].left, Slide.shapes[PicturePlaceholder[i1]].top, Slide.shapes[PicturePlaceholder[i1]].width, Slide.shapes[PicturePlaceholder[i1]].height)               
             
prs.save(RootPath+OutputName+'.pptx')

#prs = Presentation(TemplatePath)
#TypeSlideLayout = prs.slide_layouts[0]
#Slide = prs.slides.add_slide(TypeSlideLayout)
#PicturePath = PictureFilePath
#Slide.shapes.add_picture(PicturePath, Slide.shapes[2].left, Slide.shapes[2].top, Slide.shapes[2].width, Slide.shapes[2].height)
#Slide.shapes[7].name
#prs.save('D:\\基本工作資料\\my work\\PPTX\\test.pptx')
