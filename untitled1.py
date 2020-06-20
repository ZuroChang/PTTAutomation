# -*- coding: utf-8 -*-
"""
Created on Wed Jun 17 16:16:42 2020

@author: ZuroChang
"""

import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Pt


WordSize={'FontSize':[32,28,24,20,18,16,14,12,10]
             ,'Width':[415786,361553,319835,259866,237592,207893,180777,153995]
             ,'height':[633747,506997,422498,362141,362141,316873,281665,253498,162455]
             }

def GetPt(width,height,WordNum):
    WideCapacity=[width//entry for entry in WordSize['Width']]
    
    HeightCapacity=[]
    for entry in WordSize['height']:
        if entry:
            HeightCapacity.append(height//entry)
        else:
            HeightCapacity.append(None)
    
    for i0 in range(len(HeightCapacity)):
        if not HeightCapacity[len(HeightCapacity)-i0-1]:
            HeightCapacity[len(HeightCapacity)-i0-1]=HeightCapacity[len(HeightCapacity)-i0]
    
    for i0 in range(len(WordSize['FontSize'])):
        if WideCapacity[i0]*HeightCapacity[i0]>=WordNum:
            return(WordSize['FontSize'][i0])
    


PictureFolder='C:/Users/ZuroChang/PythonScript/PTTAutomation/ReportGSAM_ZuroChang_20181016/'
Import=A1
i0=0

Data=Import.iloc[i0]

Template=Data['SlideLayout']['Template']
prs = Presentation(Template+'.pptx')
Slide=prs.slides.add_slide(prs.slide_layouts[Data['SlideLayout']['SlideLayout']]) 



def InsertTitle(Shape,Title):
    Shape.text=Title
    Shape.text_frame.paragraphs[0].font.size = Pt(GetPt(Shape.width,Shape.height,len(Title)))

def InsertContent(Shape,Content):
    Shape.text=Content
    
def InsertPicture(Shape,Picture):
    Slide.shapes.add_picture(Picture,Shape.left,Shape.top,Shape.width,Shape.height)

count=1
for shape in Slide.placeholders:
    if 'Title'==shape.name[:len('Title')]:
        # print(shape.height)
        InsertTitle(shape,Data['Title'])
    elif 'Picture'==shape.name[:len('Picture')]:
        if count<=len(Data['Pictures']):
            InsertPicture(shape,PictureFolder+Data['Pictures'][count-1]+'.png')
            count+=1
    elif 'Text'==shape.name[:len('Text')]:
        text_frame=shape.text_frame
        
        p=text_frame.paragraphs[0]
        p.text=Data['Content'][0]['Text']
        p.level=Data['Content'][0]['lvl']
    
        for entry in Data['Content'][1:]:
            p=text_frame.add_paragraph()
            p.text=entry['Text']
            p.level=entry['lvl']
            
            
        # InsertContent(shape,Data['Content'])
    

prs.save('test.pptx')
